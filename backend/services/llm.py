from openai import OpenAI
from config import get_settings
import docx
from pptx import Presentation
import io

settings = get_settings()

client = OpenAI(
    api_key=settings.DEEPSEEK_API_KEY,
    base_url=settings.DEEPSEEK_BASE_URL,
)


def read_docx(file_content: bytes) -> str:
    doc = docx.Document(io.BytesIO(file_content))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def read_pptx(file_content: bytes) -> str:
    prs = Presentation(io.BytesIO(file_content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def create_docx(text: str) -> bytes:
    doc = docx.Document()
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


SIMPLE_PROMPT = """Ты — преподаватель, который объясняет сложные вещи простым языком.

Перепиши эту лекцию так, чтобы:
1. Использовать простые, понятные слова
2. Добавлять примеры из жизни
3. Структурировать информацию (подзаголовки, списки)
4. Делать акценты на главном
5. Использовать аналогии для сложных概念ов

Сохрани структуру лекции (разделы, подразделы), но перепиши содержание простым языком.

Лекция:
{content}"""

EXTENDED_PROMPT = """Ты — преподаватель, который объясняет сложные вещи простым языком и даёт дополнительную информацию.

Перепиши эту лекцию так, чтобы:
1. Использовать простые, понятные слова
2. Добавлять примеры из жизни
3. Структурировать информацию (подзаголовки, списки)
4. Делать акценты на главном
5. Использовать аналогии для сложных概念ов

После каждой секции добавляй блок "Для расширения кругозора:" с дополнительной информацией, интересными фактами, связями с другими областями знаний.

Сохрани структуру лекции (разделы, подразделы), но перепиши содержание простым языком.

Лекция:
{content}"""

SUMMARY_PROMPT = """Ты — помощник, который делает краткую выжимку из нескольких лекций.

Проанализируй эти лекции и создай:
1. Краткую выжимку (самое важное из каждой лекции)
2. Общую тему, которая объединяет лекции
3. Ключевые термины и их простые определения
4. Связи между лекциями

Формат: структурированный текст с подзаголовками.

Лекции:
{content}"""


def process_text(text: str, mode: str = "simple", extended: bool = False) -> str:
    if mode == "summary":
        prompt = SUMMARY_PROMPT.format(content=text)
    elif extended:
        prompt = EXTENDED_PROMPT.format(content=text)
    else:
        prompt = SIMPLE_PROMPT.format(content=text)

    response = client.chat.completions.create(
        model=settings.DEEPSEEK_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=4096,
    )
    return response.choices[0].message.content
